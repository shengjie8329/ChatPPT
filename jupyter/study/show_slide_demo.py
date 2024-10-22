#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
### `Slide` 幻灯片

`Slide` 表示 PowerPoint 演示文稿中的一张幻灯片。

- **属性**：
  - `shapes`：返回该幻灯片中的所有形状，包括文本框、图像、图表等。
  - `slide_layout`：返回该幻灯片的布局对象，指示使用的布局类型。

- **方法**：
  - `shapes.add_shape()`：在t 特定幻灯片上添加一个形状（例如，矩形或圆形）。

- **示例**：

```python
slide = presentation.slides.add_slide(slide_layout)  # 添加一张幻灯片
```

#### `Shape` 形状

- **概述**：`Shape` 表示幻灯片中的一个形状，可以是文本框、图片、图表、SmartArt、表格等。每个 `Shape` 对象都具有位置、大小和样式属性。

- **属性**：
  - `name`: 形状名称，对应 placeholder。
  - `left`、`top`、`width`、`height`：定义形状的位置和尺寸。
  - `text`：如果是文本框，可以访问或修改其内容。

- **方法**：
  - `add_textbox(left, top, width, height)`：用于在幻灯片上添加一个文本框。
  - `add_picture(image_path, left, top, width=None, height=None)`：用于添加图片。

- **示例**：

  ```python
  textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))  # 添加文本框
  ```

@Time   :2024/10/22 15:21
@Author :lancelot.sheng
@File   :show_slide_demo.py
"""
from pptx import Presentation
from pptx.util import Inches

presentation = Presentation("../../templates/LGBTTemplate.pptx")
slide = presentation.slides[0]  # 获取第一张幻灯片
print(slide.shapes[0].text)  # 输出幻灯片中的所有形状

# 打印每页形状 名称和文本，如果是非文本（如 PlaceholderPicture）将会报错
for idx, slide in enumerate(presentation.slides):
    print(f"[slide id]:{idx}")
    for shape in slide.shapes:
        print(f"shape name:{shape.name}")
        try:
            print(f"shape text:{shape.text}")
        except Exception as e:
            print(e)

        print("\n")

'''
### 【深入理解】Presentation 和 SlideMaster 类继承关系的 UML 类图

在这个 UML 类图中：

1. **`Presentation`** 类是顶层对象，它包含多个 `Slides` 对象。
2. **`Slides`** 类是一个幻灯片集合，通过它可以添加或访问单独的 `Slide` 对象。
3. **`Slide`** 类代表单个幻灯片，它包含形状（`Shapes`）和占位符（`SlidePlaceholders`），并且它通过布局（`SlideLayout`）来定义外观。
4. **`SlideMaster`** 类包含多个 `SlideLayouts`，它定义了幻灯片的母版布局。
5. **`SlideLayout`** 类定义了幻灯片的布局结构，其中有占位符和形状。
6. **`Shape`** 类代表幻灯片中的形状或文本框等内容。


```
+------------------------+
|      Presentation      |
+------------------------+
| - slides: Slides       |
| - slide_masters: SlideMasters |
| - slide_layouts: SlideLayouts |
| - core_properties: CoreProperties |
+------------------------+
| + save(file: str)                       |
+------------------------+
           |
           v
+--------------------+
|       Slides       |
+--------------------+
| - slides: Slide[]  |
+--------------------+
| + add_slide(layout: SlideLayout) -> Slide |
| + get(slide_id: int) -> Slide | None      |
+--------------------+
           |
           v
+--------------------+
|       Slide        |
+--------------------+
| - slide_id: int    |
| - shapes: Shapes   |
| - placeholders: SlidePlaceholders |
| - slide_layout: SlideLayout       |
+--------------------+
| + add_shape(shape: Shape)         |
| + add_picture(image: Picture)     |
| + add_table(rows: int, cols: int) |
+--------------------+
           |
           v
+--------------------+
|    SlideMaster     |
+--------------------+
| - slide_layouts: SlideLayouts[]  |
+--------------------+
| + get_by_name(name: str) -> SlideLayout |
| + index(slide_layout: SlideLayout) -> int |
+--------------------+
           |
           v
+--------------------+
|   SlideLayouts     |
+--------------------+
| - layouts: SlideLayout[] |
+--------------------+
| + remove(slide_layout: SlideLayout)      |
+--------------------+
           |
           v
+--------------------+
|   SlideLayout      |
+--------------------+
| - placeholders: SlidePlaceholders[] |
| - shapes: Shapes[]                  |
| - slide_master: SlideMaster          |
+--------------------+
           |
           v
+--------------------+
|      Shape         |
+--------------------+
| - name: str        |
| - fill: FillFormat |
| - line: LineFormat |
+--------------------+
| + add_textbox(left, top, width, height)  |
| + add_picture(image_file: str)           |
+--------------------+
```
'''

# 完整打印 Slides 每一页的所有属性
for s in presentation.slides:
    print(f"Slide ID: {s.slide_id}")
    print(f"  Layout: {s.slide_layout}")
    print(f"  Shapes: {len(s.shapes)} shapes")
    print(f"  Placeholders: {len(s.placeholders)} placeholders")

    print("  Shape Details:")
    for shape in s.shapes:
        print(f"    - Shape Name: {shape.name}, Type: {shape.shape_type}")

    print("  Placeholder Details:")
    for placeholder in s.placeholders:
        print(f"    - Placeholder ID: {placeholder.placeholder_format.idx}, Type: {placeholder.placeholder_format.type}")

    print("\n")  # Adding a new line between slides for better readability
# 输出说明：
# Slide ID: 每张幻灯片的唯一标识符。
# Layout: 使用的幻灯片布局对象。
# Shapes: 输出该幻灯片中的形状数量，并列出每个形状的详细信息（名称和类型）。
# Placeholders: 输出该幻灯片中的占位符数量，并列出每个占位符的 ID 和类型。

# 新增一页内容
# 使用 Slide_ID 获取指定页面
last_slide_layout = presentation.slides.get(slide_id=1864).slide_layout

# 新增一页幻灯片
new_slide = presentation.slides.add_slide(last_slide_layout)

# 总页数变成了 12
print(len(presentation.slides))

# 修改新增页标题
print(new_slide.shapes[0].name)
new_slide.shapes[0].text = "测试新增页面标题"
presentation.save("ChatPPT_update.pptx")

'''
## 添加页面内容的方法

### 文本`TextFrame`

- **概述**：`TextFrame` 表示一个文本框，包含文本和相关的格式设置。每个 `TextFrame` 可以包含多段文本。
- **属性**：
  - `text`：获取或设置文本框的文本内容。
  - `paragraphs`：返回文本框中的所有段落，允许对每段进行单独格式化。
- **示例**：

```python
text_frame = textbox.text_frame  # 获取文本框的文本框架
text_frame.text = "这是一段文本"  # 设置文本内容
```

#### **文本框**

在幻灯片上添加一个文本框：`add_textbox(left, top, width, height)`

```python
left = Inches(1)  # 位置
top = Inches(1)
width = Inches(5)
height = Inches(2)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame  # 获取文本框内容
text_frame.text = "这是一段文本"
```

#### **段落 `Paragraph`**

- **概述**：`Paragraph` 表示 `TextFrame` 中的单个段落。
- **属性**：
  - `text`：获取或设置段落的文本内容。
  - `font`：获取段落的字体设置，可以进行字体样式、大小和颜色的调整。
- **示例**：

```python
paragraph = text_frame.add_paragraph()  # 添加新段落
paragraph.text = "这是新的段落内容。"  # 设置段落文本
```

#### **字体 `Font`**

- **概述**：`Font` 表示字体样式，允许用户设置文本的字体样式、大小、颜色等。
- **属性**：
  - `name`：设置字体名称。
  - `size`：设置字体大小（使用 `Pt` 单位）。
  - `bold`、`italic`、`underline`：设置字体的粗体、斜体和下划线样式。
- **示例**：

```python
from pptx.util import Pt, RGBColor

run = paragraph.add_run()  # 添加文本运行
run.text = "这是加粗的文本。"
run.font.bold = True  # 设置为粗体
run.font.size = Pt(16)  # 设置字体大小
run.font.color.rgb = RGBColor(255, 0, 0)  # 设置字体颜色为红色
```
'''
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# 添加文本内容幻灯片
slide_layout = presentation.slide_layouts[-1]
slide = presentation.slides.add_slide(slide_layout)


# 打印新增页属性
print(f"Slide ID: {slide.slide_id}")
print(f"  Layout: {slide.slide_layout}")
print(f"  Shapes: {len(slide.shapes)} shapes")
print(f"  Placeholders: {len(slide.placeholders)} placeholders")

print("  Shape Details:")
for shape in slide.shapes:
    print(f"    - Shape Name: {shape.name}, Type: {shape.shape_type}")

print("  Placeholder Details:")
for placeholder in slide.placeholders:
    print(f"    - Placeholder ID: {placeholder.placeholder_format.idx}, Type: {placeholder.placeholder_format.type}")


# 填充原有布局中的占位符（标题和文本）
title = slide.shapes.title
title.text = "python-pptx 新增文本内容示例"
content = slide.placeholders[11]  #这个序号是 Placeholder ID
content.text = "填充原有的文本占位符"

# 新增文本框
left = Inches(6)
top = Inches(5)
width = Inches(5)
height = Inches(1)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = "额外的文本框内容"

# 格式化文本
paragraph = text_frame.add_paragraph()  # 添加新段落
paragraph.text = "这是一个新的段落。"  # 设置段落文本

# 设置字体
run = paragraph.add_run()  # 添加文本运行
run.text = " 这部分是加粗的文本。"  # 设置文本内容
run.font.bold = True  # 设置为粗体
run.font.size = Pt(16)  # 设置字体大小
run.font.color.rgb = RGBColor(255, 0, 0)  # 设置字体颜色为红色

# 添加图片
left = Inches(3)
top = Inches(3)
width = Inches(3)
height = Inches(3)
pic = slide.shapes.add_picture("../../images/forecast.png", left, top, width, height)


# 添加表格
rows = 2
cols = 2
left = Inches(2)
top = Inches(2)
width = Inches(4)
height = Inches(4)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = Inches(1)
table.columns[1].width = Inches(1)
table.cell(0,0).text = "1" #在指定位置写入文本
table.cell(0,1).text = "1"
table.cell(1,0).text = "1"
table.cell(1,1).text = "1"


# 保存 PPTX 文件
presentation.save("ChatPPT_append_text.pptx")