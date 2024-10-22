#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
3. Placeholder 占位符
Placeholder 是一个重要的概念，用于指示幻灯片中应该放置的内容类型（如标题、内容、图片等）。

每个布局和页面都有预定义的占位符列表（placeholders），可以用来快速添加内容。
@Time   :2024/10/22 14:15
@Author :lancelot.sheng
@File   :show_placeholder_demo.py
"""
from pptx import Presentation
from pptx.util import Inches

template = Presentation("../../templates/MasterTemplate.pptx")
slide_layouts = template.slide_layouts

slide_layout = template.slide_layouts[-2]
print(slide_layout.name)

# 打印出所有的 placeholder
for p in slide_layout.placeholders:
    print(p.name)

# LayoutPlaceholder 类型
print(type(slide_layout.placeholders[0]))

"""
### 【深入理解】Placeholder 类的继承关系（UML 类图）

在 UML 类图中，我们通常使用以下符号来表示类之间的关系：

- 正方形表示类名。
- 三个部分的矩形表示类的属性和方法。
- 箭头表示关系，如继承、关联、聚合和组合。

在这个UML类图中：

- **Placeholder** 是一个抽象的基类，包含了所有占位符共有的属性和方法。
- **MasterPlaceholder** 继承自 **Placeholder**，添加了 `name` 属性，并提供了 `add_slide` 方法。
- **LayoutPlaceholder** 继承自 **MasterPlaceholder**，添加了用于插入内容的方法，如 `insert_picture` 和 `insert_table`。
- **PicturePlaceholder** 和 **TablePlaceholder** 都是 **LayoutPlaceholder** 的子类，它们分别添加了用于设置图片和表格的特定方法，如 `set_picture` 和 `set_table`。



```
+-------------------+
|     Placeholder     |
+-------------------+
| - idx            |
| - shape_type     |
+-------------------+
| + method1()      |
| + method2()      |
+-------------------+
          ^
          |
+-------------------+
|     MasterPlaceholder |
+-------------------+
| - name            |
+-------------------+
| + add_slide()    |
+-------------------+
          ^
          |
+-------------------+
|     LayoutPlaceholder |
+-------------------+
| + insert_picture() |
| + insert_table()   |
+-------------------+
          |
          +-------------------+
          |                   |
          |                   |
          v                   v
+-------------------+     +-------------------+
|   PicturePlaceholder |     |    TablePlaceholder |
+-------------------+     +-------------------+
| + set_picture()  |     | + set_table()     |
+-------------------+     +-------------------+
```


请注意，这个类图是简化的，实际的Python-pptx库中的类可能包含更多的属性和方法。此外，UML类图通常包含更多的细节，如可见性（如 `+` 表示公共，`-` 表示私有）和关系类型（如泛化、实现、关联、依赖等）。在这个简化的示例中，只展示了泛化（继承）关系，并且所有的属性和方法都被假设为公共的。
"""

