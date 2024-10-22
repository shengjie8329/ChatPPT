#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
python-pptx 库抽象与母版中的概念对应
在 python-pptx 中，可以将一些抽象概念与 PowerPoint 母版中的内容和布局对应起来：

母版（Master Slide）：

在 python-pptx 中，可以通过 presentation.slide_master 来访问母版。母版包含了幻灯片的基本格式和样式，可以定义统一的外观。
布局（Layouts）：

使用 presentation.slide_layouts 可以访问不同的幻灯片布局，例如标题幻灯片、内容幻灯片等。每种布局都有预定义的占位符，可以用于快速插入内容。
内容类型（Content Types）：

对应于文本框、图片、图表等内容类型，可以使用 add_textbox()、add_picture()、add_table()、add_chart() 等方法来添加这些元素。

### 1. `Presentation` 类

- **概述**：

`Presentation` 是 `python-pptx` 中的核心类，用于表示一个 PowerPoint 演示文稿。

- **构造方法**：

创建一个新的空演示文稿。
@Time   :2024/10/22 13:56
@Author :lancelot.sheng
@File   :cteate_presentation_demo.py
"""
from pptx import Presentation
from pptx.util import Inches

# 创建一个新的 PowerPoint 文件
presentation = Presentation()

# 保存 PPTX 文件
presentation.save("empty_presentation.pptx")